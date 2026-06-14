import os
from io import BytesIO

# Import extractors with fallback error handling if packages aren't installed yet
try:
    from pypdf import PdfReader
except ImportError:
    PdfReader = None

try:
    import docx
except ImportError:
    docx = None

try:
    import pptx
except ImportError:
    pptx = None


def extract_text_from_pdf(file_bytes: bytes) -> str:
    """Extracts text from PDF bytes."""
    if PdfReader is None:
        raise ImportError("El paquete 'pypdf' no está instalado. Instálalo con 'pip install pypdf'.")
    
    pdf_file = BytesIO(file_bytes)
    reader = PdfReader(pdf_file)
    text_content = []
    
    for i, page in enumerate(reader.pages):
        text = page.extract_text()
        if text:
            text_content.append(f"--- Página {i+1} ---\n{text}")
            
    return "\n\n".join(text_content)


def extract_text_from_docx(file_bytes: bytes) -> str:
    """Extracts text from Word (.docx) bytes."""
    if docx is None:
        raise ImportError("El paquete 'python-docx' no está instalado. Instálalo con 'pip install python-docx'.")
    
    docx_file = BytesIO(file_bytes)
    doc = docx.Document(docx_file)
    text_content = []
    
    for i, paragraph in enumerate(doc.paragraphs):
        if paragraph.text.strip():
            text_content.append(paragraph.text)
            
    # Also extract text from tables
    for table in doc.tables:
        for row in table.rows:
            row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if row_text:
                text_content.append(" | ".join(row_text))
                
    return "\n".join(text_content)


def extract_text_from_pptx(file_bytes: bytes) -> str:
    """Extracts text from PowerPoint (.pptx) bytes."""
    if pptx is None:
        raise ImportError("El paquete 'python-pptx' no está instalado. Instálalo con 'pip install python-pptx'.")
    
    pptx_file = BytesIO(file_bytes)
    prs = pptx.Presentation(pptx_file)
    text_content = []
    
    for i, slide in enumerate(prs.slides):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())
        if slide_text:
            text_content.append(f"--- Diapositiva {i+1} ---\n" + "\n".join(slide_text))
            
    return "\n\n".join(text_content)


def extract_text(file_bytes: bytes, file_name: str) -> str:
    """Determines file type by name extension and extracts its text."""
    ext = os.path.splitext(file_name)[1].lower()
    
    if ext == ".pdf":
        return extract_text_from_pdf(file_bytes)
    elif ext == ".docx":
        return extract_text_from_docx(file_bytes)
    elif ext in [".pptx", ".ppt"]:
        return extract_text_from_pptx(file_bytes)
    else:
        raise ValueError(f"Formato de archivo no soportado: {ext}. Utilice PDF, DOCX o PPTX.")
